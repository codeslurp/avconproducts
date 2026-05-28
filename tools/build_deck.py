"""Build an AVCON-branded PowerPoint deck describing the Product Code Finder.

Output: valve-selector/AVCON_Product_Code_Finder_Overview.pptx (sibling of
README.txt / MAINTENANCE.txt so it's easy to find).

Audience: AVCON Controls leadership / management. Focus on business value
(speed, error reduction, attach rate), with screenshots to make the workflow
concrete. Less detail on implementation than an engineering pitch would have.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ----- Paths -----------------------------------------------------------------
HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
SCREENS = ROOT / "screenshots"
LOGO = ROOT / "app" / "static" / "avcon-logo.png"
OUT_PATH = ROOT / "AVCON_Product_Code_Finder_Overview.pptx"

# ----- Brand palette ---------------------------------------------------------
TEAL          = RGBColor(0x01, 0x7E, 0x80)
TEAL_DARK     = RGBColor(0x01, 0x66, 0x68)
TEAL_SOFT_BG  = RGBColor(0xE6, 0xF1, 0xF1)
NEAR_BLACK    = RGBColor(0x3A, 0x3A, 0x3A)
SLATE         = RGBColor(0x4B, 0x4F, 0x58)
GRAY_TEXT     = RGBColor(0x7A, 0x7A, 0x7A)
LIGHT_GRAY    = RGBColor(0xE5, 0xE5, 0xE5)
OFF_WHITE     = RGBColor(0xFC, 0xFC, 0xFC)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
WARNING       = RGBColor(0xC9, 0x7A, 0x00)

FONT_SANS  = "Calibri"           # broadest availability; Roboto is a free download
FONT_SERIF = "Cambria"           # stand-in for Roboto Slab — both serif, both common
FONT_MONO  = "Consolas"

# ----- Geometry --------------------------------------------------------------
# 16:9 widescreen (PowerPoint default)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_L = Inches(0.6)
MARGIN_R = Inches(0.6)
MARGIN_T = Inches(0.5)
MARGIN_B = Inches(0.5)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R


# ----- Helpers ---------------------------------------------------------------
def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def add_rect(slide, x, y, w, h, fill, line_fill=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_fill is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_fill
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *,
             font=FONT_SANS, size=18, bold=False, italic=False,
             color=NEAR_BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, x, y, w, h, paragraphs, *,
                   font=FONT_SANS, size=14, color=SLATE,
                   line_spacing=1.25, bullet_size=None):
    """paragraphs: list of either str OR (text, bold) OR dict with keys."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, item in enumerate(paragraphs):
        if isinstance(item, dict):
            text = item.get("text", "")
            cfg = item
        elif isinstance(item, tuple):
            text, bold = item
            cfg = {"text": text, "bold": bold}
        else:
            text = item
            cfg = {"text": text}

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = cfg.get("align", PP_ALIGN.LEFT)
        p.line_spacing = cfg.get("line_spacing", line_spacing)
        if i > 0:
            p.space_before = Pt(cfg.get("space_before", 6))
        run = p.add_run()
        run.text = text
        run.font.name = cfg.get("font", font)
        run.font.size = Pt(cfg.get("size", size))
        run.font.bold = cfg.get("bold", False)
        run.font.italic = cfg.get("italic", False)
        run.font.color.rgb = cfg.get("color", color)
    return tb


def add_header_bar(slide, slide_num, total, section=""):
    """Top thin teal accent bar + slide number bottom-right."""
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.1), TEAL)
    add_text(slide, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.35),
             Inches(0.6), Inches(0.25),
             f"{slide_num} / {total}",
             font=FONT_SANS, size=9, color=GRAY_TEXT, align=PP_ALIGN.RIGHT)
    if section:
        add_text(slide, MARGIN_L, SLIDE_H - Inches(0.35),
                 Inches(6.0), Inches(0.25),
                 section.upper(),
                 font=FONT_SANS, size=9, bold=True, color=TEAL)


def add_logo(slide, x, y, height):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), x, y, height=height)


def add_title(slide, title, subtitle=None, *, y=None):
    if y is None:
        y = Inches(0.45)
    add_text(slide, MARGIN_L, y, CONTENT_W, Inches(0.7),
             title, font=FONT_SERIF, size=32, bold=True, color=NEAR_BLACK)
    if subtitle:
        add_text(slide, MARGIN_L, y + Inches(0.65), CONTENT_W, Inches(0.4),
                 subtitle, font=FONT_SANS, size=14, italic=True, color=GRAY_TEXT)
    add_rect(slide, MARGIN_L, y + Inches(1.1), Inches(0.6), Inches(0.06), TEAL)


def add_screenshot(slide, name, *, x, y, w, h):
    """Insert screenshot scaled to fit a w x h box, preserving aspect ratio.
    Adds a subtle gray border around it."""
    path = SCREENS / name
    if not path.exists():
        # placeholder rectangle
        add_rect(slide, x, y, w, h, LIGHT_GRAY)
        add_text(slide, x, y, w, h, f"[missing: {name}]",
                 size=12, color=GRAY_TEXT, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        return
    # We let python-pptx scale via height to preserve aspect ratio, then
    # check whether width exceeds the target box; if so, scale by width.
    pic = slide.shapes.add_picture(str(path), x, y, width=w)
    if pic.height > h:
        slide.shapes._spTree.remove(pic._element)
        pic = slide.shapes.add_picture(str(path), x, y, height=h)
        pic.left = x + int((w - pic.width) / 2)
    else:
        pic.top = y + int((h - pic.height) / 2)
        pic.left = x + int((w - pic.width) / 2)
    # subtle border
    pic.line.color.rgb = LIGHT_GRAY
    pic.line.width = Pt(0.75)


# ----- Slides ----------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 12

    # ------- Slide 1: Title -------
    s = add_blank_slide(prs)
    add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, OFF_WHITE)
    # left accent panel
    add_rect(s, Emu(0), Emu(0), Inches(0.4), SLIDE_H, TEAL)
    # logo top-right corner
    add_logo(s, SLIDE_W - Inches(2.2), Inches(0.5), Inches(0.9))
    # super-title
    add_text(s, Inches(1.0), Inches(2.4), Inches(11), Inches(0.5),
             "AVCON CONTROLS PVT. LTD.",
             font=FONT_SANS, size=16, bold=True, color=TEAL)
    # main title
    add_text(s, Inches(1.0), Inches(2.95), Inches(11), Inches(1.4),
             "Product Code Finder",
             font=FONT_SERIF, size=54, bold=True, color=NEAR_BLACK)
    # subtitle
    add_text(s, Inches(1.0), Inches(4.3), Inches(11), Inches(0.6),
             "From hours of engineering work to seconds of sales clicks.",
             font=FONT_SERIF, size=20, italic=True, color=SLATE)
    # spec line
    add_text(s, Inches(1.0), Inches(5.1), Inches(11), Inches(0.4),
             "49,289 SKUs · 5 product families · One desktop tool",
             font=FONT_SANS, size=14, color=GRAY_TEXT)
    # divider
    add_rect(s, Inches(1.0), Inches(5.7), Inches(2.0), Inches(0.04), TEAL)
    add_text(s, Inches(1.0), Inches(5.85), Inches(11), Inches(0.4),
             "Internal overview · For AVCON Leadership",
             font=FONT_SANS, size=12, color=GRAY_TEXT)

    # ------- Slide 2: The Problem -------
    s = add_blank_slide(prs)
    add_header_bar(s, 2, TOTAL, section="The Problem")
    add_title(s, "Quoting a valve + actuator is engineering work, done by sales.")
    add_paragraphs(s, MARGIN_L, Inches(2.1), Inches(12), Inches(4.5), [
        ("Today, a quote for a customer who wants 'a 2-inch CF8M ball valve, actuated' takes a salesperson through:", False),
        {"text": "    1.  Open the master Excel data sheets. Pick a valve series, then size, then body and ball materials, seat, characteristics, end connection, and ball type.", "size": 13},
        {"text": "    2.  Compute the valve's break torque (BTO) and apply the safety factor: FOS = BTO × 1.5.", "size": 13},
        {"text": "    3.  Open the actuator data sheet. Pick a model with enough torque at the customer's supply pressure.", "size": 13},
        {"text": "    4.  Cross-check that the actuator's mounting (PCD, shaft) matches the valve's stem.", "size": 13},
        {"text": "    5.  Repeat for the spring-return variant if the customer wants fail-safe behavior.", "size": 13},
        ("", False),
        ("The result is slow, inconsistent across salespeople, and dependent on engineering judgement at the point of sale — every quote risks an under-sized actuator or a mounting mismatch.", True),
    ], size=14, color=SLATE)

    # ------- Slide 3: What we built -------
    s = add_blank_slide(prs)
    add_header_bar(s, 3, TOTAL, section="Solution")
    add_title(s, "A desktop tool that replaces the data sheets.",
              subtitle="Open a browser, pick attributes from dropdowns, get the SKU. No spreadsheets, no manual sizing.")
    add_screenshot(s, "01-home-empty.png",
                   x=Inches(1.5), y=Inches(2.3), w=Inches(10.3), h=Inches(4.5))
    add_text(s, MARGIN_L, SLIDE_H - Inches(0.65), CONTENT_W, Inches(0.25),
             "Landing page — AVCON-branded, all 5 product families behind two pickers.",
             font=FONT_SANS, size=11, italic=True, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

    # ------- Slide 4: How it works — step 1, pick a category -------
    s = add_blank_slide(prs)
    add_header_bar(s, 4, TOTAL, section="Workflow · 1 of 4")
    add_title(s, "Step 1 — Pick the product family.")
    add_paragraphs(s, MARGIN_L, Inches(1.85), Inches(5.8), Inches(4.5), [
        ("Two pickers, side by side:", True),
        ("  Valves — Ball Valve, Butterfly Valve (Centric)", False),
        ("  Actuators — Pneumatic (Rack & Pinion, Scotch Yoke), Electrical (Rotary)", False),
        ("", False),
        ("Each option shows its SKU count and source data file — so the salesperson and the leadership reviewer always know where the numbers come from.", False),
        ("", False),
        ("The salesperson can keep both pickers open simultaneously to compare a valve quote with its matching actuator.", False),
    ], size=13)
    add_screenshot(s, "02-valves-picker-open.png",
                   x=Inches(7.1), y=Inches(1.85), w=Inches(5.8), h=Inches(5.0))

    # ------- Slide 5: Step 2, the cascade -------
    s = add_blank_slide(prs)
    add_header_bar(s, 5, TOTAL, section="Workflow · 2 of 4")
    add_title(s, "Step 2 — Walk the cascade. The app blocks invalid combinations.")
    add_paragraphs(s, MARGIN_L, Inches(1.85), Inches(5.8), Inches(4.8), [
        ("After picking a family (here: Ball Valve), the salesperson configures the SKU through dropdowns:", False),
        ("", False),
        ("  Series · Size · Body Material · Ball Material · Seat Material · Characteristics · End Connection · Ball Type", True),
        ("", False),
        ("Each dropdown only shows options that lead to a real SKU. If the customer asks for a combination AVCON does not manufacture, the dropdowns simply do not offer it — no invalid quotes are possible.", False),
        ("", False),
        ("Series and Model fields support type-ahead search for power users.", False),
    ], size=13)
    add_screenshot(s, "04-ball-valve-empty.png",
                   x=Inches(7.1), y=Inches(1.85), w=Inches(5.8), h=Inches(5.0))

    # ------- Slide 6: Step 3, result -------
    s = add_blank_slide(prs)
    add_header_bar(s, 6, TOTAL, section="Workflow · 3 of 4")
    add_title(s, "Step 3 — The SKU resolves with everything needed for the quote.")
    add_paragraphs(s, MARGIN_L, Inches(1.85), Inches(5.8), Inches(4.8), [
        ("Once the cascade is complete, the Result panel shows the four numbers a quote needs:", False),
        ("", False),
        ("  Bare Valve Code — internal SKU", False),
        ("  Catalogue Code — customer-facing reference", False),
        ("  BTO (N·m) — break torque", False),
        ("  FOS = BTO × 1.5 — pre-computed safety-factored torque target", False),
        ("", False),
        ("'All attributes' expands the full 50-field engineering data sheet for that SKU — design standard, leakage class, certifications, weights, stem dimensions, everything.", False),
    ], size=13)
    add_screenshot(s, "05-ball-valve-resolved.png",
                   x=Inches(7.1), y=Inches(1.85), w=Inches(5.8), h=Inches(5.0))

    # ------- Slide 7: Step 4, the recommendation -------
    s = add_blank_slide(prs)
    add_header_bar(s, 7, TOTAL, section="Workflow · 4 of 4")
    add_title(s, "Step 4 — The tool recommends the matching actuator.")
    # left column: explanation
    add_paragraphs(s, MARGIN_L, Inches(1.85), Inches(5.8), Inches(5.0), [
        ("This is the feature that removes engineering judgement from the sales path.", True),
        ("", False),
        ("Every ball-valve SKU in the catalog already names a compatible Rack & Pinion actuator model. The app reads that pairing and surfaces it as a one-click button under the Result panel.", False),
        ("", False),
        ("Clicking it opens the Actuator section with the model pre-selected. The salesperson finishes the 4-5 customer-driven choices (body material, certification, etc.) — they never have to size torque or check mounting compatibility.", False),
        ("", False),
        ("Both sections stay visible side by side, so the valve quote and the actuator quote can be assembled together.", False),
    ], size=13)
    add_screenshot(s, "06-actuator-pre-selected.png",
                   x=Inches(7.1), y=Inches(1.5), w=Inches(5.8), h=Inches(5.6))

    # ------- Slide 8: Business impact -------
    s = add_blank_slide(prs)
    add_header_bar(s, 8, TOTAL, section="Business Impact")
    add_title(s, "What changes for AVCON.")
    # Big comparison table-ish layout
    table_x = MARGIN_L
    table_y = Inches(2.1)
    col_w = Inches(4.0)
    row_h = Inches(0.95)
    gap = Inches(0.15)

    headers = ["Metric", "Today (manual)", "With the Code Finder"]
    rows = [
        ("Time per quote (valve + actuator)",
         "Minutes — open spreadsheets, size torque, cross-check mounting.",
         "Seconds — cascade, click 'Configure matching actuator', done."),
        ("Engineering knowledge required",
         "High — salesperson does BTO × 1.5 and reads torque-vs-pressure tables.",
         "None — the catalog's engineering-approved pairing is used directly."),
        ("Consistency across the sales team",
         "Two salespeople may quote different actuators for the same valve.",
         "Same valve → same recommended actuator, every time."),
        ("Attach rate (actuator with valve)",
         "Customer often comes back later for the actuator — or buys elsewhere.",
         "Every valve quote naturally includes the actuator suggestion."),
        ("Risk of an under-sized actuator",
         "Real — depends on salesperson reading torque tables correctly.",
         "Eliminated for standard pairings — engineering signed off in the catalog."),
    ]

    # Column headers
    for i, h_text in enumerate(headers):
        x = table_x + (col_w + gap) * i
        add_rect(s, x, table_y, col_w, Inches(0.4), TEAL)
        add_text(s, x + Inches(0.15), table_y + Inches(0.07), col_w - Inches(0.3), Inches(0.3),
                 h_text, font=FONT_SANS, size=12, bold=True, color=WHITE)

    # Data rows
    for r_i, row in enumerate(rows):
        ry = table_y + Inches(0.4) + Inches(0.08) + row_h * r_i
        for c_i, cell in enumerate(row):
            x = table_x + (col_w + gap) * c_i
            fill = TEAL_SOFT_BG if c_i == 2 else (LIGHT_GRAY if c_i == 1 else WHITE)
            line = LIGHT_GRAY
            add_rect(s, x, ry, col_w, row_h - Inches(0.05), fill, line_fill=line)
            color = NEAR_BLACK if c_i == 0 else (SLATE if c_i == 1 else TEAL_DARK)
            bold = c_i == 0
            size = 10 if c_i > 0 else 11
            add_text(s, x + Inches(0.15), ry + Inches(0.1), col_w - Inches(0.3), row_h - Inches(0.2),
                     cell, font=FONT_SANS, size=size, bold=bold, color=color)

    # ------- Slide 9: Catalog coverage -------
    s = add_blank_slide(prs)
    add_header_bar(s, 9, TOTAL, section="Coverage")
    add_title(s, "What's in the tool today.",
              subtitle="All five product families load on startup. 49,289 SKUs total.")
    # 5 cards in a row
    families = [
        ("Ball Valve",            "3,892",  "valve SKUs",  "Valves"),
        ("Butterfly Valve",       "43,200", "valve SKUs",  "Valves"),
        ("Rack & Pinion",         "1,620",  "actuator SKUs", "Actuators — Pneumatic"),
        ("Scotch Yoke",           "481",    "actuator SKUs", "Actuators — Pneumatic"),
        ("Electrical Rotary",     "96",     "actuator SKUs", "Actuators — Electrical"),
    ]
    card_w = Inches(2.35)
    card_h = Inches(2.1)
    gap = Inches(0.10)
    total_w = card_w * 5 + gap * 4
    start_x = (SLIDE_W - total_w) / 2
    card_y = Inches(2.3)
    for i, (name, count, unit, group) in enumerate(families):
        x = start_x + (card_w + gap) * i
        add_rect(s, x, card_y, card_w, card_h, WHITE, line_fill=LIGHT_GRAY)
        add_rect(s, x, card_y, card_w, Inches(0.08), TEAL)
        add_text(s, x + Inches(0.2), card_y + Inches(0.25), card_w - Inches(0.4), Inches(0.4),
                 group.upper(),
                 font=FONT_SANS, size=9, bold=True, color=TEAL)
        add_text(s, x + Inches(0.2), card_y + Inches(0.6), card_w - Inches(0.4), Inches(0.6),
                 name,
                 font=FONT_SERIF, size=18, bold=True, color=NEAR_BLACK)
        add_text(s, x + Inches(0.2), card_y + Inches(1.25), card_w - Inches(0.4), Inches(0.5),
                 count,
                 font=FONT_SERIF, size=24, bold=True, color=TEAL_DARK)
        add_text(s, x + Inches(0.2), card_y + Inches(1.7), card_w - Inches(0.4), Inches(0.3),
                 unit,
                 font=FONT_SANS, size=10, italic=True, color=GRAY_TEXT)

    # context line below cards
    add_paragraphs(s, MARGIN_L, Inches(5.0), CONTENT_W, Inches(1.5), [
        ("All catalogs are loaded from the master AVCON Excel data sheets — the same documents engineering already maintains. No second source of truth is created.", False),
        ("", False),
        ("New product families can be added by dropping a new Excel file in the data folder and registering it with a single config block in the app.", False),
    ], size=13, color=SLATE)

    # ------- Slide 10: How it stays current -------
    s = add_blank_slide(prs)
    add_header_bar(s, 10, TOTAL, section="Maintenance")
    add_title(s, "Engineering updates the Excel. The tool follows.",
              subtitle="No developer in the loop for routine catalog changes.")
    add_paragraphs(s, MARGIN_L, Inches(2.1), CONTENT_W, Inches(5), [
        ("Adding a SKU, fixing a value, or populating more actuator pairings:", True),
        ("    1.  Open the relevant Excel file from the data folder. Edit. Save.", False),
        ("    2.  Restart the application (close the command window, relaunch).", False),
        ("    3.  Hard-refresh the browser. Done.", False),
        ("", False),
        ("Changes are read once at startup. The app picks the most recently modified file matching each product family, so old versions can be kept as one-step rollback.", False),
        ("", False),
        ("Structural changes (renaming sheets, reordering columns) require a small config edit. A separate MAINTENANCE.txt document in the application folder walks the engineering team through both routine and structural updates.", False),
        ("", False),
        ("The application is local-only. No data leaves the salesperson's machine — important for catalog confidentiality.", True),
    ], size=13, color=SLATE)

    # ------- Slide 11: What's next -------
    s = add_blank_slide(prs)
    add_header_bar(s, 11, TOTAL, section="Roadmap")
    add_title(s, "What's next.")
    items = [
        ("Populate the remaining actuator pairings",
         "Today only one recommendation cell per ball-valve row is filled (col 49). The catalog already has space for up to 9 alternatives — Double Acting variants, Spring-Return Fail-Close, Spring-Return Fail-Open. Populating these in Excel unlocks an even richer recommendation list with no further development work."),
        ("Extend pairings to Scotch Yoke and Electrical Rotary",
         "Current pairings only resolve to Rack & Pinion actuators. Adding model references for the other two actuator families in the same way will let the tool recommend across the full actuator catalog."),
        ("Multi-user deployment",
         "Today the app runs on individual salesperson laptops. A single shared instance (intranet) would mean engineering's Excel updates propagate to everyone the moment they restart — no need to redistribute files."),
        ("PDF / Print export",
         "Browsers can already print the result panel; a dedicated 'Generate quote PDF' button would produce a customer-ready data sheet in one click."),
    ]
    iy = Inches(2.0)
    for i, (head, body) in enumerate(items):
        y = iy + Inches(1.18) * i
        # bullet square
        add_rect(s, MARGIN_L, y + Inches(0.18), Inches(0.12), Inches(0.12), TEAL)
        add_text(s, MARGIN_L + Inches(0.3), y, CONTENT_W - Inches(0.3), Inches(0.4),
                 head, font=FONT_SERIF, size=15, bold=True, color=NEAR_BLACK)
        add_text(s, MARGIN_L + Inches(0.3), y + Inches(0.42),
                 CONTENT_W - Inches(0.3), Inches(0.75),
                 body, font=FONT_SANS, size=11, color=SLATE)

    # ------- Slide 12: Closing -------
    s = add_blank_slide(prs)
    add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, OFF_WHITE)
    add_rect(s, Emu(0), Emu(0), Inches(0.4), SLIDE_H, TEAL)
    add_logo(s, SLIDE_W - Inches(2.2), Inches(0.5), Inches(0.9))
    add_text(s, Inches(1.0), Inches(2.5), Inches(11), Inches(1.2),
             "Thank you.",
             font=FONT_SERIF, size=54, bold=True, color=NEAR_BLACK)
    add_text(s, Inches(1.0), Inches(3.75), Inches(11), Inches(0.6),
             "AVCON Product Code Finder — ready for the sales team to use today.",
             font=FONT_SERIF, size=20, italic=True, color=SLATE)
    add_rect(s, Inches(1.0), Inches(4.6), Inches(2.0), Inches(0.04), TEAL)
    add_text(s, Inches(1.0), Inches(4.8), Inches(11), Inches(0.4),
             "Questions, feedback, or requests for the next iteration: please reach out.",
             font=FONT_SANS, size=14, color=GRAY_TEXT)

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")
    print(f"Size:  {OUT_PATH.stat().st_size:,} bytes")


if __name__ == "__main__":
    build()
